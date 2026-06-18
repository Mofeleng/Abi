import os
import matplotlib
matplotlib.use('Agg')  
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def _generate_chart_image(visual_meta: dict, output_path: str) -> str:
    chart_type = visual_meta.get("chart_type", "bar")
    title = visual_meta.get("title", "Data Insight")
    labels = visual_meta.get("labels", [])
    values = visual_meta.get("values", [])
    
    if not labels or not values:
        return None
        
    plt.figure(figsize=(6, 4.5))
    colors = ['#1e3a8a', '#3b82f6', '#10b981', '#f59e0b', '#ef4444', '#8b5cf6']
    
    if chart_type == "pie":
        plt.pie(values, labels=labels, autopct='%1.1f%%', colors=colors[:len(labels)], startangle=140)
    elif chart_type == "line":
        plt.plot(labels, values, marker='o', color='#3b82f6', linewidth=2)
        plt.grid(True, linestyle='--', alpha=0.5)
    else:  
        plt.bar(labels, values, color=colors[:len(labels)], edgecolor='none')
        plt.grid(axis='y', linestyle='--', alpha=0.3)
        
    plt.title(title, fontsize=12, fontweight='bold', pad=15, color='#1e293b')
    plt.xticks(rotation=20, fontsize=8, ha='right')
    plt.yticks(fontsize=9)
    plt.tight_layout()
    
    sanitized_title = "".join([c for c in title if c.isalnum() or c in (' ', '_')]).rstrip().replace(' ', '_').lower()
    img_path = output_path.replace(".pptx", f"_{sanitized_title}.png")
    plt.savefig(img_path, dpi=200, transparent=True)
    plt.close()
    return img_path

def run(slide_outline: dict, output_path: str) -> str:
    print("[Abi Design] Generating a flawless 10/10 executive presentation...")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)  
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    for slide_data in slide_outline.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        
        # ─── RENDERING THE COVER PAGE LAYOUT ──────────────────────────
        if slide_data.get("type") == "cover_page":
            # Giant Centered Header Box
            cover_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(3.0))
            tf = cover_box.text_frame
            tf.word_wrap = True
            
            p1 = tf.paragraphs[0]
            p1.text = slide_data.get("title", "Supermarket Intelligence Report")
            p1.font.size = Pt(44)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A) # Dark Slate
            p1.space_after = Pt(20)
            
            p2 = tf.add_paragraph()
            p2.text = slide_data.get("subtitle", "Automated Multi-Agent Analytics Framework")
            p2.font.size = Pt(20)
            p2.font.color.rgb = RGBColor(0x4B, 0x55, 0x63) # Gray Accent
            continue
            
        # ─── RENDERING STANDARD DATA ANALYSIS SLIDES ──────────────────
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1.0))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "Executive Metric Slide")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
        
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.0), Inches(5.0))
        ctf = content_box.text_frame
        ctf.word_wrap = True
        
        for bullet in slide_data.get("bullet_points", []):
            bp = ctf.add_paragraph()
            bp.text = f"• {bullet}"
            bp.font.size = Pt(15)
            bp.space_after = Pt(12)
            bp.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
            
        visual_meta = slide_data.get("visual_element")
        if visual_meta and visual_meta.get("type") == "chart":
            try:
                chart_img = _generate_chart_image(visual_meta, output_path)
                if chart_img and os.path.exists(chart_img):
                    slide.shapes.add_picture(chart_img, Inches(7.0), Inches(1.8), width=Inches(5.8))
            except Exception as chart_err:
                print(f"[Abi Design] Chart visualization layer skipped: {chart_err}")

    prs.save(output_path)
    print(f"[Abi Design] Success. 10/10 Masterpiece saved to {output_path}")
    return output_path